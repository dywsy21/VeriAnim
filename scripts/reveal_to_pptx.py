#!/usr/bin/env python3
"""Convert a Reveal.js HTML deck to a PPTX while preserving animated GIFs.

The converter renders each slide in Chromium and uses that screenshot as the
slide background. Animated GIF elements are then re-inserted as original GIF
files at their measured browser positions, so PowerPoint keeps them animated
instead of flattening them into the screenshot.
"""

from __future__ import annotations

import argparse
import asyncio
from pathlib import Path
from urllib.parse import unquote, urlparse

from pptx import Presentation
from pptx.util import Inches
from playwright.async_api import async_playwright


SLIDE_WIDTH_PX = 1920
SLIDE_HEIGHT_PX = 1080
PPT_WIDTH_IN = 13.3333333333
PPT_HEIGHT_IN = 7.5
PX_PER_IN = SLIDE_WIDTH_PX / PPT_WIDTH_IN


def path_from_file_uri(uri: str) -> Path | None:
    parsed = urlparse(uri)
    if parsed.scheme != "file":
        return None
    if parsed.netloc:
        return Path(unquote(f"//{parsed.netloc}{parsed.path}"))
    raw_path = unquote(parsed.path)
    if len(raw_path) >= 3 and raw_path[0] == "/" and raw_path[2] == ":":
        raw_path = raw_path[1:]
    return Path(raw_path)


def px_to_inches(value: float) -> float:
    return value / PX_PER_IN


async def capture_slides(html_path: Path, assets_dir: Path) -> list[dict]:
    assets_dir.mkdir(parents=True, exist_ok=True)
    html_uri = html_path.resolve().as_uri()
    slide_records: list[dict] = []

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(
            viewport={"width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX},
            device_scale_factor=1,
        )
        await page.goto(html_uri, wait_until="networkidle")
        await page.wait_for_function("window.Reveal && Reveal.isReady && Reveal.isReady()", timeout=30_000)
        await page.evaluate(
            """() => {
                Reveal.configure({ controls: false, progress: false });
                document.documentElement.style.background = "white";
                document.body.style.background = "white";
                Reveal.layout();
            }"""
        )
        slide_count = await page.evaluate("Reveal.getSlides().length")

        for index in range(slide_count):
            await page.evaluate("(i) => { Reveal.slide(i, 0, 0); Reveal.layout(); }", index)
            await page.wait_for_timeout(350)
            screenshot_path = assets_dir / f"slide_{index + 1:02d}.png"
            await page.screenshot(path=str(screenshot_path), full_page=False)
            gifs = await page.evaluate(
                """() => {
                    const current = Reveal.getCurrentSlide();
                    return Array.from(current.querySelectorAll("img"))
                        .map((img) => {
                            const rect = img.getBoundingClientRect();
                            return {
                                src: img.currentSrc || img.src,
                                left: rect.left,
                                top: rect.top,
                                width: rect.width,
                                height: rect.height,
                                visible: rect.width > 0 && rect.height > 0
                            };
                        })
                        .filter((item) => item.visible && item.src.split("?")[0].toLowerCase().endsWith(".gif"));
                }"""
            )
            slide_records.append({"screenshot": screenshot_path, "gifs": gifs})

        await browser.close()

    return slide_records


def build_pptx(slide_records: list[dict], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(PPT_WIDTH_IN)
    prs.slide_height = Inches(PPT_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]

    for record in slide_records:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(record["screenshot"]),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        for gif in record["gifs"]:
            gif_path = path_from_file_uri(gif["src"])
            if gif_path is None or not gif_path.exists():
                continue
            slide.shapes.add_picture(
                str(gif_path),
                Inches(px_to_inches(gif["left"])),
                Inches(px_to_inches(gif["top"])),
                width=Inches(px_to_inches(gif["width"])),
                height=Inches(px_to_inches(gif["height"])),
            )

    if len(prs.slides) and len(prs.slides._sldIdLst) > len(slide_records):
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


async def main_async(args: argparse.Namespace) -> None:
    html_path = args.html.resolve()
    assets_dir = args.assets.resolve()
    output_path = args.output.resolve()
    slide_records = await capture_slides(html_path, assets_dir)
    build_pptx(slide_records, output_path)
    gif_count = sum(len(record["gifs"]) for record in slide_records)
    print(f"Wrote {output_path}")
    print(f"Slides: {len(slide_records)}")
    print(f"Animated GIF overlays: {gif_count}")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert a Reveal.js HTML deck to PPTX with animated GIF overlays.")
    parser.add_argument("html", type=Path, help="Path to Reveal.js index.html")
    parser.add_argument("-o", "--output", type=Path, required=True, help="Output PPTX path")
    parser.add_argument("--assets", type=Path, default=Path("slides/pptx_assets"), help="Directory for rendered slide backgrounds")
    return parser.parse_args()


def main() -> None:
    asyncio.run(main_async(parse_args()))


if __name__ == "__main__":
    main()
